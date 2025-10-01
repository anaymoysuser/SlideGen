from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import json
from pptx import Presentation

from pptx.enum.text import PP_ALIGN
from pathlib import Path
from utils.wei_utils import  style_bullet_content 
COLOR_BG     = RGBColor(255, 255, 255)   # 背景白
COLOR_BORDER = RGBColor(11,39,128)    # 主题蓝
COLOR_TITLE  = RGBColor(26, 26, 26)      # 深灰
COLOR_TEXT   = RGBColor(51, 51, 51)

BORDER_W = Pt(15)                         # 边框粗细
PAGE_MARGIN = Inches(2)                # 内边距


def add_paragraph(text_frame, text, size=32, bold=False, align=PP_ALIGN.CENTER):
    p = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = align


def add_cover_slide(prs, title_blocks, authors_line, date_line):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    draw_soft_background(slide, sw, sh)

    # ---- 标题文本框 ----
    title_box = slide.shapes.add_textbox(
        PAGE_MARGIN, Inches(1.2),
        sw - 2*PAGE_MARGIN, Inches(2.5)
    )
    title_box.text_frame.word_wrap = True
    fill_textbox(title_box, title_blocks)      

    # ---- 作者 / 日期 ----
    meta_box = slide.shapes.add_textbox(
        PAGE_MARGIN, Inches(4.2),
        sw - 2*PAGE_MARGIN, Inches(1.0)
    )
    tf = meta_box.text_frame
    add_paragraph(tf, authors_line, size=20, align=PP_ALIGN.CENTER)
    add_paragraph(tf, date_line,    size=18, align=PP_ALIGN.CENTER)



from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

COLOR_BLUE = RGBColor(0, 51, 153)     
COLOR_GRAY = RGBColor(245, 245, 245)   

# ----- Layout constants -----
PAGE_MARGIN = Inches(0.8)
def _apply_blue(run):
    """Apply deep-blue theme color, & handle rPr creation safely."""
    run.font.fill.solid(); run.font.fill.fore_color.rgb = COLOR_BLUE


def _apply_blue_font(font, size_pt: int, bold: bool = True):
    """Set font properties (size, bold) and apply deep‑blue color in a robust way."""
    font.size = Pt(size_pt)
    font.bold = bold
    font.fill.solid()
    font.fill.fore_color.rgb = COLOR_BLUE


def _set_run_size(run, size_pt: int):
    run.font.size = Pt(size_pt)

def add_toc_slide(prs, section_names):
    """Create a TOC slide matching the reference layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # === fixed dimensions ===
    grey_top = PAGE_MARGIN
    grey_left = PAGE_MARGIN
    box_w = sw - PAGE_MARGIN * 2
    box_h = Cm(13)  # fixed height for grey box
    bar_h = Inches(0.8) * 3  # fixed thick bar
    bar_top = (sh - bar_h) / 2  # centered blue bar (independent)
    # === independent blue bar centered in slide ===
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0, bar_top,
        sw, bar_h,
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_BLUE
    bar.line.fill.background()


    # === grey rounded rectangle ===
    grey = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        grey_left, grey_top,
        box_w, box_h,
    )
    grey.fill.solid(); grey.fill.fore_color.rgb = COLOR_GRAY
    grey.line.fill.background()

    # === vertical "CONTENTS" (8 separate lines) with tighter spacing ===
    vert_tb = slide.shapes.add_textbox(
        grey_left + box_w - Inches(1.2),
        grey_top + Inches(0.5),
        Inches(1), box_h - Inches(1.0)
    )
    vtf = vert_tb.text_frame
    vtf.clear()
    for idx, ch in enumerate("CONTENTS"):
        para = vtf.add_paragraph() if idx else vtf.paragraphs[0]
        para.text = ch
        para.alignment = PP_ALIGN.CENTER
        para.line_spacing = Pt(40)
        _apply_blue_font(para.font, 40)

    # === section list ===
    list_tb = slide.shapes.add_textbox(
        grey_left + Inches(1.0), grey_top + Inches(1.0),
        box_w - Inches(2.5), box_h - Inches(1.5)
    )
    ltf = list_tb.text_frame
    ltf.clear(); ltf.word_wrap = True

    for i, name in enumerate(section_names, 1):
        para = ltf.add_paragraph() if ltf.text else ltf.paragraphs[0]
        para.text = f"{i}. {name}"
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = Pt(42)
        _apply_blue_font(para.font, 38)




def draw_soft_background(slide, sw, sh):
    """纯白底 + 蓝色四边框""" 
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
 
    edges = [
        (0, 0, sw, BORDER_W),                     # top
        (0, sh - BORDER_W, sw, BORDER_W),         # bottom
        (0, 0, BORDER_W, sh),                     # left
        (sw - BORDER_W, 0, BORDER_W, sh),         # right
    ]
    for left, top, width, height in edges:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BORDER
        shape.line.fill.background()  
ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

def fill_textbox(textbox, blocks):
    """
    blocks = [ {alignment, bullet, level, font_size, runs:[{text,bold}]} , ... ]
    """
    tf = textbox.text_frame
    tf.clear()
    for i, blk in enumerate(blocks):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.alignment = ALIGN_MAP.get(blk.get("alignment", "left"), PP_ALIGN.LEFT)
        p.level = blk.get("level", 0)
        p.font.size = Pt(blk.get("font_size", 24))

        if blk.get("bullet"):
            p.text = ""  
            for run_data in blk["runs"]:
                r = p.add_run()
                r.text = run_data["text"]
                r.font.size = Pt(blk["font_size"])
                if run_data.get("bold"):
                    r.font.bold = True
                if r.font.color is None:
                    _ = r.font.color
              
        else:
        
            r = p.add_run()
            r.text = "".join(run["text"] for run in blk["runs"])
            r.font.size = Pt(blk["font_size"])
            blk["runs"][0]['color'] = COLOR_BORDER
            if blk["runs"][0].get("bold"):
                r.font.bold = True
            if r.font.color is None:
                _ = r.font.color
          






def ppt_from_json(json_path, ppt_path, page_size=(16, 9)):
 
    slides_data = json.load(open(json_path, "r"))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(page_size[0]), Inches(page_size[1])
    blank = prs.slide_layouts[6]

    for page_idx, page in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        sw, sh = prs.slide_width, prs.slide_height
        draw_soft_background(slide, sw, sh)

        # ---- Title ----
        if "title" in page:
            title_box = slide.shapes.add_textbox(
                PAGE_MARGIN, Inches(0.8), sw - 2*PAGE_MARGIN, Inches(1.5)
            )
            title_box.text_frame.word_wrap = True
            title_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            fill_textbox(title_box, page["title"])

        # ---- Main textbox1 ----
        if "textbox1" in page:
            tb1 = slide.shapes.add_textbox(
                PAGE_MARGIN, Inches(2.3), sw - 2*PAGE_MARGIN, sh - Inches(3)
            )
            tb1.text_frame.word_wrap = True
            tb1.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            fill_textbox(tb1, page["textbox1"])

      
        for key, y_offset in [("textbox2", 5.5), ("textbox3", 7.2)]:
            if key in page:
                tb = slide.shapes.add_textbox(
                    PAGE_MARGIN, Inches(y_offset),
                    sw - 2*PAGE_MARGIN, Inches(1.5)
                )
                tb.text_frame.word_wrap = True
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                fill_textbox(tb, page[key])

    prs.save(ppt_path)
    print("Saved:", ppt_path)


def generate_multi_slide_ppt(
    title_arrangement,
    figure_arrangement,
    text_arrangement,
    content,
    slide_width,
    slide_height,
    save_path,
    theme=None,
):
     
    prs = Presentation()
    prs.slide_width  = slide_width   
    prs.slide_height = slide_height
    blank = prs.slide_layouts[6]



   
    cover_title_blocks = content[0]["title"]
    authors_line = " ".join(run["text"] for run in content[0]["textbox1"][0]["runs"])
    date_line    = "2025/7/12"                      
    add_cover_slide(prs, cover_title_blocks, authors_line, date_line)
 
    for section in content:
        style_bullet_content(section["title"], COLOR_BORDER, COLOR_BORDER)
    
    
    def extract_title_text(title_blocks): 
        lines = []
        for block in title_blocks:
            line = "".join(run["text"] for run in block.get("runs", []))
            lines.append(line)
        return " ".join(lines)
 
    section_names = [extract_title_text(section["title"]) for section in content[1:]]

    add_toc_slide(prs, section_names)

 

    for idx, section in enumerate(content):
        if idx==0:
            continue
        slide = prs.slides.add_slide(blank)
        draw_soft_background(slide, prs.slide_width, prs.slide_height)

        # ========== 1. Title ==========
        tbox_cfg = title_arrangement[idx]
        t_shape = slide.shapes.add_textbox(
            Inches(tbox_cfg["x"]), Inches(tbox_cfg["y"]),
            Inches(tbox_cfg["width"]), Inches(tbox_cfg["height"])
        )
        t_shape.text_frame.word_wrap = True
        
        fill_textbox(t_shape, section["title"])

        # ========== 2. Figure   ==========
        fbox_cfg = figure_arrangement[idx]
        fig_path = fbox_cfg.get("figure_path")
     
        if fig_path and Path(fig_path).is_file() and fbox_cfg["width"] > 0.1:
            slide.shapes.add_picture(
                fig_path,
                Inches(fbox_cfg["x"]), Inches(fbox_cfg["y"]),
                width  = Inches(fbox_cfg["width"]),
                height = Inches(fbox_cfg["height"]),
            )

        # ========== 3. Textbox1 ==========
        tb_cfg = text_arrangement[idx]
        text_shape = slide.shapes.add_textbox(
            Inches(tb_cfg["x"]), Inches(tb_cfg["y"]),
            Inches(tb_cfg["width"]), Inches(tb_cfg["height"])
        )
        text_shape.text_frame.word_wrap = True
        fill_textbox(text_shape, section["textbox1"])

 
        for key in ("textbox2", "textbox3"):
            if key in section:
             
                y_offset = tb_cfg["y"] + tb_cfg["height"] + 0.2  # 0.2 inch 间隔
                extra = slide.shapes.add_textbox(
                    Inches(tb_cfg["x"]), Inches(y_offset),
                    Inches(tb_cfg["width"]), Inches(1.2)
                )
                extra.text_frame.word_wrap = True
                extra.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                fill_textbox(extra, section[key])

    prs.save(save_path)
    print(" PPT saved →", save_path)



if __name__ == "__main__":
    ppt_from_json(
        json_path="step_slides.json",
        ppt_path="step_soft_template.pptx",
        page_size=(13.33, 7.5)  # 16:9  
    )
