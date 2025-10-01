
import json
from pathlib import Path
from typing import Dict, List
import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.dml.color import MSO_THEME_COLOR
 
 
 
from pprint import pprint

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE   
from pptx.util import Inches

import json
import difflib
from pptx import Presentation

COLOR_WHITE = RGBColor(0, 0, 0) 
THEME_COLOR = RGBColor(0, 0, 0) 
 

def _insert_picture_keep_ratio(ph, img_path: Path):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image

    slide = ph.part.slide   
    ph_left, ph_top = ph.left, ph.top
    ph_width, ph_height = ph.width, ph.height
 
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih
 
    if ph_width / ph_height > aspect:
        new_h = ph_height
        new_w = int(new_h * aspect)
    else:
        new_w = ph_width
        new_h = int(new_w / aspect)
 
    new_left = ph_left + int((ph_width - new_w) / 2)
    new_top = ph_top + int((ph_height - new_h) / 2)
 
    pic = slide.shapes.add_picture(str(img_path), new_left, new_top, width=new_w, height=new_h)
 
    # --- move picture *behind* all text placeholders ---
    spTree = slide.shapes._spTree
    spTree.remove(pic.element)         # temporarily take it out
    spTree.insert(2, pic.element)      # index 0=background,1=layout,2≈最底自定义层

    # --- finally, remove the now‑unused picture placeholder itself ---
    ph.element.getparent().remove(ph.element)
  
from pptx.util import Pt
from PIL import Image

def insert_image_below_content(slide, img_path: Path):
    """
    Insert an image below the lowest existing shape (text or image) on the slide.
    Centered horizontally. Resizes if not enough space.
    """ 
    with Image.open(img_path) as img:
        width_px, height_px = img.size
    aspect_ratio = width_px / height_px
 
    slide_width = slide.part.slide_layout.part.package.presentation_part.slide_width
    slide_height = slide.part.slide_layout.part.package.presentation_part.slide_height
 
    target_width = slide_width * 0.6
    target_height = target_width / aspect_ratio
 
    lowest_bottom = 0
    for shape in slide.shapes:
        bottom = shape.top + shape.height
        if bottom > lowest_bottom:
            lowest_bottom = bottom

    margin = Pt(20)
    available_space = slide_height - lowest_bottom - margin

    if available_space < target_height:
        target_height = available_space
        target_width = target_height * aspect_ratio
        if target_height <= 0:
            print("Not enough space to insert image:", img_path)
            return

    
    left = (slide_width - target_width) // 2
    top = lowest_bottom + margin

    slide.shapes.add_picture(str(img_path), left, top, width=int(target_width), height=int(target_height))

 

TEXT_TYPES = {
    PH_TYPE.TITLE,
    PH_TYPE.CENTER_TITLE,
    PH_TYPE.SUBTITLE,
    PH_TYPE.BODY,
}

def find_text_placeholders(slide): 
    """Return (part_num_ph, subsection_ph, body_ph) by position."""
    txt_ph = [
        s for s in slide.shapes
        if (
            s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
            s.has_text_frame and
            s.placeholder_format.type in TEXT_TYPES   
        )
    ]

    if len(txt_ph) < 3:
        raise ValueError("Not enough text placeholders on this slide")

     
    txt_ph.sort(key=lambda s: s.top)
 
    first_row = sorted(txt_ph[:2], key=lambda s: s.left)
    part_ph, title_ph = first_row
    body_ph = txt_ph[2]          

    return part_ph, title_ph, body_ph

def get_content(sec_title, sub_title,outline):
    _, sub = _best_match(outline, sec_title, sub_title)
    return sub.get("content", "") if sub else ""

def set_font_color(paragraph, theme_color):
    paragraph.font.fill.solid()
    if theme_color is not None:
        paragraph.font.fill.fore_color.theme_color = theme_color
    else:
        paragraph.font.fill.fore_color.rgb = RGBColor(255, 105, 180)  # pink fallback

def insert_visuals_auto(slide, visuals: list[Path]):
    """
    Automatically insert visuals (images, tables, formulas) into all available
    picture placeholders on the given slide.
    """
    # Find all picture placeholder shapes
    picture_placeholders = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
           "Picture" in shape.name
    ]

    if len(visuals) > len(picture_placeholders):
        print(f"Warning: not enough picture placeholders on slide (needed {len(visuals)}, found {len(picture_placeholders)})")
        # remaining = visuals[len(picture_placeholders):]
        # for img_path in remaining:
        #     insert_image_below_content(slide, Path(img_path))

    #  Insert images one by one
    for img_path, ph in zip(visuals, picture_placeholders):
        _insert_picture_keep_ratio(ph, Path(img_path))


def _placeholder_by_name(slide, name: str):
    """Return placeholder shape whose .name == name."""

    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f'Placeholder "{name}" not found on slide master.')
 
 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def scan_layout_placeholders(template_path: str):
    prs = Presentation(template_path)
    layout_map = {}

    for layout in prs.slide_layouts:
        layout_name = layout.name
        placeholder_names = []

        for shape in layout.shapes:
            shape_type = shape.shape_type
            shape_type_name = str(shape_type)
            try:
                shape_type_name = MSO_SHAPE_TYPE(shape_type).name
            except ValueError:
                pass

            placeholder_names.append({
                "name": shape.name,
                "type": shape_type_name
            })

        layout_map[layout_name] = placeholder_names

    return layout_map


from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_theme_color_from_title(prs, layout_index=2):
    slide_layout = prs.slide_layouts[layout_index]
    temp_slide = prs.slides.add_slide(slide_layout)
 
    for shape in temp_slide.shapes:
        if (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and 
            shape.placeholder_format.type == 1  # 1 == TITLE
        ):
            para = shape.text_frame.paragraphs[0]
            font_color = para.font.color
            if font_color.type == 2:  # 2 == THEME
                theme_color = font_color.theme_color 
                xml_slides = prs.slides._sldIdLst  
                xml_slides.remove(xml_slides[-1])  # remove last slide
                return theme_color
    return None
 

def resolve_visual_paths(slide_info, args):
    from pathlib import Path
    paper = args.paper_name
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
    base_dir = os.path.join(prefix, paper)   
 
    images_json  = json.load(open(os.path.join( prefix, f"{paper}_images.json")))
    tables_json  = json.load(open(os.path.join( prefix, f"{paper}_tables.json")))
    if  args.formula_mode == 3:
        formulas_dir = os.path.join("contents", paper, "formula_images")
    else: 
        formulas_dir = base_dir
 
    def _match_img_key(key: str, mapping: dict):
        if key in mapping:
            return mapping[key]
        stem = Path(key).stem
        if stem in mapping:
            return mapping[stem] 
        no_pref = stem.split("_", 1)[-1]
        if no_pref in mapping:
            return mapping[no_pref]
        for k in mapping:
            if Path(k).stem == stem or Path(k).stem == no_pref:
                return mapping[k]

        return None
    def _inspect_images_json(images_json, limit=12):
        print("\n[inspect] images_json type:", type(images_json).__name__)
        if isinstance(images_json, dict):
            keys = list(images_json.keys())
            print(f"[inspect] total keys: {len(keys)}; sample keys:", keys[:limit])
            if keys:
                k0 = keys[0]
                print("[inspect] sample value for first key:")
                pprint(images_json[k0])
        elif isinstance(images_json, list):
            print(f"[inspect] list length: {len(images_json)}; sample items (first {limit}):")
            pprint(images_json[:limit])
        else:
            print("[inspect] unexpected images_json structure")

    def _resolve_and_check(path1: Path):
        if not path1.exists():
            raise FileNotFoundError(f"Missing visual file: {path1}")
        return path1

    # -------- images --------
    image_paths = []
    for img_str in slide_info["images"]:
        name = Path(str(img_str)).name  # 只取文件名，避免目录中的数字干扰
        m = re.search(r'(?:picture|image|fig|table|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        if m:
            img_id = m.group(1)
        else:
            # 回退：取文件名里“最后一段数字”
            nums = re.findall(r'\d+', name)
            if not nums:
                raise ValueError(f"Invalid image ID format: {img_str}")
            img_id = nums[-1]

        img_id = str(int(img_id))  # 归一化（去前导零）
        print("img_str:", img_str, "  file:", name, "  img_id:", img_id)

        # _inspect_images_json(images_json)
        rec = _match_img_key(img_id, images_json)
        if rec is None:
            from pathlib import Path
            cand = {img_id, Path(img_id).stem, Path(img_id).stem.split("_", 1)[-1]}
            print("[resolve_visual_paths] cannot match:", img_id, "tried:", cand)
            print("[resolve_visual_paths] example keys:", list(images_json.keys())[:8])
            print(sorted(images_json.keys()))
            raise KeyError(f"Image id not found: {img_id}")
        target = rec['image_path']
        print("target: ",target) 
        image_paths.append(_resolve_and_check( Path(target)))

    # -------- tables --------
    table_paths = []
    for tb_str in slide_info["tables"]:

        name = Path(str(tb_str)).name  
        m = re.search(r'(?:picture|image|fig|table|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        if m:
            tb_id = m.group(1)
        else:
            # 回退：取文件名里“最后一段数字”
            nums = re.findall(r'\d+', name)
            if not nums:
                raise ValueError(f"Invalid image ID format: {tb_str}")
            tb_id = nums[-1]

        tb_id = str(int(tb_id))  # 归一化（去前导零）
        print("[Debug] tb_str:", tb_str, "  file:", name, "  img_id:", tb_id)
  
        print(f"[Debug] tables_json.keys() = {list(tables_json.keys())[:10]}")  # 打印前10个 key

        match = _match_img_key(tb_id, tables_json)
        if match is None:
            raise ValueError(f"[resolve_visual_paths] Cannot find table match for tb_id: {tb_id}")
        target = match['table_path']

        # target = _match_img_key(tb_id, tables_json)['table_path']
        if target is None:
            raise ValueError(f"Table ID '{tb_id}' not found in mapping.")
        image_paths.append(_resolve_and_check( Path(target)))

    # -------- formulas --------
    formula_paths = []
    for fname in slide_info["formulas"]:
        if args.formula_mode == 3:        
            final_path = formulas_dir / fname
        else:
            final_path = resolve_formula_mode1_path(fname,args)
        formula_paths.append(_resolve_and_check(final_path))

    return image_paths + table_paths + formula_paths



_TMPL_TO_PH = {
    "T1_TextOnly":      [],
    "T2_ImageRight":    ["pic_right"],
    "T3_ImageLeft":     ["pic_left"],
    "T4_ImageTop":      ["pic_top"],
    "T5_TwoImages":     ["pic_left", "pic_right"],
   
    # Grid-based layouts
    "T7_2x2_TopImage":    ["img1", "img2"],              # 2x2: top two blocks are images, bottom two are text
    "T8_2x2_BottomImage": ["img3", "img4"],              # 2x2: bottom two blocks are images, top two are text
    "T9_2x2_AltTextImg":  ["img1", "img3"],              # 2x2: alt image and text diagonally
    "T10_4Img_2x2Grid":   ["img1", "img2", "img3", "img4"],  # New: 2x2 image-only layout
 
    # Three images + one text (split by variant)
    "T11_3Img_TopTextBottom": ["img1", "img2", "img3"],  # Top 3 images (horizontal), text below
    "T12_3Img_BottomTextTop": ["img1", "img2", "img3"],  # Top text, bottom 3 square images in a row
    "T13_3Img":   ["img1", "img2", "img3"],      # Title on top, three images below

}


 


  
import re
from pathlib import Path
import difflib

 
def _extract_idx(val):
    """从 int / '7' / 'image_7.png' / 'table7' 中提取编号 7；失败返回 None。"""
    if isinstance(val, int):
        return val
    m = re.findall(r"\d+", str(val))
    return int(m[-1]) if m else None

def _nums_from_files(files):
    """['image_7.png', 'table_1.png'] -> {7, 1}"""
    out = set()
    for f in files:
        m = re.findall(r"\d+", str(f))
        if m:
            out.add(int(m[-1]))
    return out

def _best_match(data, sec_title: str, sub_title: str, min_ratio: float = 0.55):
   
    best_sec, best_sec_score = None, 0.0
    for sec in data.get("sections", []):
        s = difflib.SequenceMatcher(None, sec.get("title","").lower(), (sec_title or "").lower()).ratio()
        if s > best_sec_score:
            best_sec, best_sec_score = sec, s
    if not best_sec or best_sec_score < min_ratio:
        return None, None

    best_sub, best_sub_score = None, 0.0
    for sub in best_sec.get("subsections", []):
        s = difflib.SequenceMatcher(None, sub.get("title","").lower(), (sub_title or "").lower()).ratio()
        if s > best_sub_score:
            best_sub, best_sub_score = sub, s
    if not best_sub or best_sub_score < min_ratio:
        return None, None
    return best_sec, best_sub
 
def _collect_reasons_for_kind(sec, sub, files, figs_data, kind: str):
   
    _, sub_d = _best_match(figs_data, sec, sub)
    if not sub_d:
        return ""

    want = _nums_from_files(files)
    if not want:
        return ""
 
    pairs = []  # [(N, value_idx)]
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith(kind):  
            continue
        N = _extract_idx(k)   
        v_idx = _extract_idx(v) 
        if N is None or v_idx is None:
            continue
        if v_idx in want:
            pairs.append((N, v_idx))

    reasons = []
    
    only_one_asset = (len(pairs) == 1)

    for N, _v_idx in pairs:
       
        candidate_keys = [
            f"reason{N}",
        ]
        if kind == "table":
            candidate_keys.append(f"reasonT{N}")
        else:
            candidate_keys.extend([f"reasonI{N}", f"reasonImg{N}"])
     
        if only_one_asset:
            candidate_keys.append("reason")
 
        found = None
        for rk in candidate_keys:
            if rk in sub_d and isinstance(sub_d[rk], str) and sub_d[rk].strip():
                found = sub_d[rk].strip()
                break
 
        if not found and "reason" in sub_d and isinstance(sub_d["reason"], str) and sub_d["reason"].strip():
            found = sub_d["reason"].strip()

        if found:
            reasons.append(found)

    return "\n".join(reasons)

def get_image_reasons(sec, sub, image_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, image_files, figs_data, kind="image")

def get_table_reasons(sec, sub, table_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, table_files, figs_data, kind="table")

def get_formula_reasons(sec, sub, formula_files, formula_data) -> str:
    
    _, sub_d = _best_match(formula_data, sec, sub)
    if not sub_d:
        return ""
 
    pairs = []
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith("formula"):
            continue
        N = _extract_idx(k)   
        if N is None:
            continue
        rkey = f"reason{N}"
        rtxt = sub_d.get(rkey)
        if isinstance(rtxt, str) and rtxt.strip():
            pairs.append((N, rtxt.strip()))
 
    try:
        pairs.sort(key=lambda x: int(x[0]))
    except Exception:
        
        pass

     
    if not pairs and isinstance(sub_d.get("reason"), str) and sub_d["reason"].strip():
        if formula_files and len(formula_files) == 1:
            return sub_d["reason"].strip()
 
    if formula_files and len(pairs) > len(formula_files):
        pairs = pairs[:len(formula_files)]
 
    out, seen = [], set()
    for _, r in pairs:
        if r not in seen:
            out.append(r)
            seen.add(r)

    return "\n".join(out)

def resolve_formula_mode1_path(fname: str, args) -> Path:
    """
    Extract formula index i from fname (like "formula_4.png"),
    and generate the path:
    <{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-i.png
    """
    match = re.search(r'formula[_-](\d+)', fname)
    if not match:
        raise ValueError(f"Cannot extract index from formula filename: {fname}")
    
    i = match.group(1)
    path_str = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-{i}.png"
    return Path(path_str)

import json
from pathlib import Path
from typing import Dict, Any, List

def _is_T1_textonly(s: Dict[str, Any]) -> bool:
  
    return (
        s.get("template_id") == "T1_TextOnly"
        and not s.get("images")
        and not s.get("tables")
        and not s.get("formulas")
    )

def pair_T1_to_T19(plan_path: str, write_back: bool = True) -> int:
     
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    out: List[Dict[str, Any]] = []

    i, n, made = 0, len(slides), 0
    while i < n:
        cur = slides[i]
        if (
            i + 1 < n
            and _is_T1_textonly(cur)
            and _is_T1_textonly(slides[i + 1])
            and cur.get("section") == slides[i + 1].get("section")
        ):
            left, right = cur, slides[i + 1]
            out.append({
                "section": cur.get("section"),
                "template_id": "T19_2Text",
                "columns": [
                    {
                        "subsection": left.get("subsection", "") or "",
                        "bullets": left.get("bullets", []) or []
                    },
                    {
                        "subsection": right.get("subsection", "") or "",
                        "bullets": right.get("bullets", []) or []
                    }
                ]
            })
            made += 1
            i += 2
            continue
 
        out.append(cur)
        i += 1

    plan["slides"] = out
    if write_back:
        p.write_text(json.dumps(plan, ensure_ascii=False, indent=4), encoding="utf-8")
    return made

def validate_no_consecutive_T1(plan_path: str) -> List[int]:
 
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    bad_idxs: List[int] = []
    for i in range(len(slides) - 1):
        a, b = slides[i], slides[i + 1]
        if _is_T1_textonly(a) and _is_T1_textonly(b) and a.get("section") == b.get("section"):
            bad_idxs.append(i)
    return bad_idxs
from pptx.util import Pt

def _clear_text_frame(tf):
    if tf.paragraphs:
        tf.paragraphs[0].text = ""
  
        while len(tf.paragraphs) > 1:
            p = tf._element.p_lst[-1]
            p.getparent().remove(p)
    else:
        tf.clear()

def _fill_bullets(tf, bullets, lvl0_size=24, lvl1_size=24):
    _clear_text_frame(tf)
    for b in (bullets or []):
       
        p = tf.add_paragraph()
        p.text = (b.get("text") or "").strip()
        p.level = 0
        p.font.size = Pt(lvl0_size)
       
        for s in (b.get("sub") or []):
            sp = tf.add_paragraph()
            sp.text = str(s).strip()
            sp.level = 1
            sp.font.size = Pt(lvl1_size)
from pptx.enum.shapes import PP_PLACEHOLDER

def _get_placeholder(slide, name):
   
    for shp in slide.shapes:
        if getattr(shp, "name", "").strip().lower() == name.strip().lower():
            return shp
    return None

def _ph_text_n(slide, n:int):
    
    targets = {f"text placeholder {n}", f"文本占位符 {n}"}
    for shp in slide.shapes:
        nm = getattr(shp, "name", "").strip().lower()
        if nm in targets:
            return shp
         
        if ("placeholder" in nm  ) and nm.endswith(f" {n}"):
            return shp
    return None
def _ph_by_idx(slide, idx:int):
    for ph in getattr(slide, "placeholders", []):
        if ph.placeholder_format.idx == idx:
            return ph
    return None

def fill_T19_2Text(slide, slide_info, section_no_text):
     
    print("slide_info:")
    print(slide_info)
    part_ph = (
        _get_placeholder(slide, "Part")
        or _get_placeholder(slide, "Text Placeholder 2")
        or _ph_text_n(slide, 2)
    )
   
    title_bar = _get_placeholder(slide, "Text Placeholder 1") or _ph_text_n(slide, 1)
    title_bar = _ph_by_idx(slide, 1)  or _ph_text_n(slide, 1)
 
    section_title = slide_info.get("section") 
    
    title_bar.text = section_title
    tf = title_bar.text_frame
    tf.clear()                        
    tf.paragraphs[0].text = section_title
    print("[AFTER] title_bar text =", repr(tf.text))

    lt = (
        _get_placeholder(slide, "Left Title")
        or _get_placeholder(slide, "Text Placeholder 3")
        or _ph_text_n(slide, 3)
    )
    lb = (
        _get_placeholder(slide, "Left Body")
        or _get_placeholder(slide, "Text Placeholder 4")
        or _ph_text_n(slide, 4)
    )
    rt = (
        _get_placeholder(slide, "Right Title")
        or _get_placeholder(slide, "Text Placeholder 5")
        or _ph_text_n(slide, 5)
    )
    rb = (
        _get_placeholder(slide, "Right Body")
        or _get_placeholder(slide, "Text Placeholder 6")
        or _ph_text_n(slide, 6)
    )
 
    cols = slide_info.get("columns") or []
    left  = cols[0] if len(cols) > 0 else {}
    right = cols[1] if len(cols) > 1 else {}
 
    if part_ph is not None and getattr(part_ph, "has_text_frame", False):
        part_ph.text_frame.text = f"{section_no_text}"

    if title_bar is not None and getattr(title_bar, "has_text_frame", False):
        title_txt = slide_info.get("section", "") or slide_info.get("title", "")
        title_bar.text_frame.text = title_txt

    if lt is not None and getattr(lt, "has_text_frame", False):
        lt.text_frame.text = left.get("subsection", "") or left.get("title", "") or ""

    if rt is not None and getattr(rt, "has_text_frame", False):
        rt.text_frame.text = right.get("subsection", "") or right.get("title", "") or ""

    if lb is not None and getattr(lb, "has_text_frame", False):
        _fill_bullets(lb.text_frame, left.get("bullets"))

    if rb is not None and getattr(rb, "has_text_frame", False):
        _fill_bullets(rb.text_frame, right.get("bullets"))
 
    missing = []
    for name, ph in [
        ("Part(2)", part_ph), ("Left Title(3)", lt), ("Left Body(4)", lb),
        ("Right Title(5)", rt), ("Right Body(6)", rb),
    ]:
        if ph is None:
            missing.append(name)
    if missing:
        print(f"[WARN] T19_2Text 模板缺少占位：{', '.join(missing)}")



   

def generate_pptx_from_plan( 
    args,
    template: Path | int 
):
 
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
 
    figs_json_path  =  f"contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_figures.json"
    formula_json_path = f"contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_formula_match.json"
    paper_outline_json = f'contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_raw_content.json' 
    with open(paper_outline_json, "r", encoding="utf-8") as f: outline_json  = json.load(f)
    with open(figs_json_path, encoding="utf-8") as f: figs_data   = json.load(f)
    with open(formula_json_path, encoding="utf-8") as f: formula_data   = json.load(f)
    
    plan_json = f'contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_slide_plan.json'
     
    made = pair_T1_to_T19(plan_json)    
    print(f"[plan] T1->T19 pairs made: {made}")
   
    plan: Dict = json.loads(Path(plan_json).read_text(encoding="utf-8"))


    title    = outline_json["metadata"]["title"]
    subtitle = outline_json["metadata"]["author"]
 
    # template_path = f"utils/slides_template/slides{template}_template.pptx"
    template_path = f"utils/slides_template/slides{template}_template.pptx"
    prs = Presentation(template_path)
      
    theme_color = extract_theme_color_from_title(prs)
    print("Theme color:", theme_color)
     
    # test 
    print("Available slide layouts in template:")
    for layout in prs.slide_layouts:
        print("-", layout.name)
 
    # ---------- cover ----------
    cover_layout = prs.slide_layouts.get_by_name("Title Slide")
    cover = prs.slides.add_slide(cover_layout)

    _placeholder_by_name(cover, "Title 1").text = title
    _placeholder_by_name(cover, "Subtitle 2").text = subtitle

    # ---------- Contents ----------
    outline_layout = prs.slide_layouts.get_by_name("Mulu")
    outline = prs.slides.add_slide(outline_layout)

    tf = _placeholder_by_name(outline, "Text Placeholder 1").text_frame
    tf.clear()
    seen = set() 
    unique_sections = []
 
    for slide in plan["slides"]:
        sec = slide["section"]
        if sec not in seen:
            seen.add(sec)
            unique_sections.append(sec)
 
    for i, sec in enumerate(unique_sections):
        if i == 0:
            p = tf.paragraphs[0]  # Use first paragraph to avoid empty line
        else:
            p = tf.add_paragraph()
        p.text = sec
        p.level = 0
        p.font.size = Pt(36) 
        # p.font.fill.solid()
        # set_font_color(p, theme_color)
            
    # ---------- Body ----------
    current_section  = None
    section_counter  = 0
    for slide_info in plan["slides"]: 
        # ---------- Content ----------
        if slide_info["section"] != current_section:
            current_section = slide_info["section"]
            section_counter += 1 
            section_layout = prs.slide_layouts.get_by_name("dan_mulu")
            sec_slide = prs.slides.add_slide(section_layout)
            for shape in sec_slide.shapes:
                print(f"Shape: {shape.name}")
            _placeholder_by_name(sec_slide, "Text Placeholder 2").text = f"PART {section_counter:02d}"
            _placeholder_by_name(sec_slide, "Title 1").text = current_section
  
        template_id = slide_info["template_id"]
        layout = prs.slide_layouts.get_by_name(template_id)
          
        if layout is None:
            raise ValueError(f" Template layout '{template_id}' not found in template.")
        slide = prs.slides.add_slide(layout)
    
        if template_id == "T19_2Text":
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.is_placeholder:
                    if shape.has_text_frame:
                        print(f" Name: {shape.name}")
                        print(f"  Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}")
                        print(f"  Text: '{shape.text_frame.text.strip()}'")
    
            fill_T19_2Text(slide, slide_info, section_no_text=f"{section_counter:02d}")
            continue

        part_ph, title_ph, body_ph = find_text_placeholders(slide)
  
        part_ph.text = f"{section_counter:02d}" 
        title_ph.text = slide_info["subsection"]
  
   
        # bullets + sub-bullets
        if body_ph:
            tf = body_ph.text_frame
            if tf.paragraphs:
                tf.paragraphs[0].text = ""   
            else:
                tf.clear() 
            for bullet in slide_info["bullets"]:
                p = tf.add_paragraph()
                p.text, p.level = bullet["text"], 0
                p.font.size = Pt(24)
                for sub in bullet.get("sub", []):
                    sp = tf.add_paragraph()
                    sp.text, sp.level = sub, 1
                    sp.font.size = Pt(24)
         
         
        visuals = resolve_visual_paths(slide_info, args)
        insert_visuals_auto(slide, visuals)
 
 
        notes_chunks = [] 
        txt = get_content(slide_info["section"], slide_info["subsection"], outline_json)
        if txt: notes_chunks.append(txt)
        if slide_info.get("images"):
            img_r = get_image_reasons(slide_info["section"], slide_info["subsection"],
                                      slide_info["images"], figs_data)
            notes_chunks.append(img_r)
        if slide_info.get("tables"):
            tb_r = get_table_reasons(slide_info["section"], slide_info["subsection"],
                                     slide_info["tables"], figs_data)
            notes_chunks.append(tb_r)
        if slide_info.get("formulas"):
            fm_r = get_formula_reasons(
                slide_info["section"],
                slide_info["subsection"],
                slide_info["formulas"],
                formula_data  ,
                )
                
            notes_chunks.append(fm_r)
                    
        if notes_chunks:
            nframe = slide.notes_slide.notes_text_frame
            if nframe.text and not nframe.text.endswith("\n"):
                nframe.text += "\n"
            nframe.text += "\n\n".join(notes_chunks)

 
    thanks_layout = prs.slide_layouts.get_by_name("Last_page")
    thanks = prs.slides.add_slide(thanks_layout)
    title_ph = _placeholder_by_name(thanks, "Title 1")
    title_ph.text = "THANKS!"
 
    run = title_ph.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    output_pptx = f'contents/{args.paper_name}/{args.model_name_t}_{args.model_name_v}_output_slides.pptx' 
    prs.save(str(output_pptx))
       
 
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Fill PPTX from slides plan.")
    parser.add_argument(
        "--plan",  
        default="/home/ubuntu/yifan/eeg2fmri/Paper2PPTX/contents/STEP_A_General_and_Scalable_Framework_for_Solving_Video_Inverse_Problems/<4o_4o>_slide_plan.json",  
        help="slides_plan.json"
        )
    parser.add_argument(
        "--paper_name",  
        default="STEP_A_General_and_Scalable_Framework_for_Solving_Video_Inverse_Problems"  
        )
    parser.add_argument(
        "--template",  
        type=int,
        default=3,
        help="Template number, e.g. 3 for slides3_template.pptx"
    )
    parser.add_argument(
        "--out", 
        default="output.pptx" 
        )
    parser.add_argument(
        "--model_name_t", 
        default="4o" 
        )
    parser.add_argument(
        "--model_name_v", 
        default="4o" 
        )
        
    args = parser.parse_args()

    layout_info = scan_layout_placeholders("/home/ubuntu/yifan/eeg2fmri/Paper2PPTX/utils/slides_template/slides3_template.pptx")

    for layout_name, shapes in layout_info.items():
        print(f"\n Layout: {layout_name}")
        for s in shapes:
            print(f"  - {s['name']} ({s['type']})")

 

    generate_pptx_from_plan(args.plan, args.template, args.out)
     
    print(f" Saved to {args.out}")
