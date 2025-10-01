# -*- coding: utf-8 -*-
"""
Render LaTeX formulas from a JSON file into PNG images and insert them into a PPT.

Input JSON format example (array of objects):
[
  {
    "page": 2,
    "bbox": {"x": 123.4, "y": 567.8, "w": 90.1, "h": 22.3},
    "latex_raw": "\\mathcal { L } _ { d i s t i l l } = ...",
    "latex": "\\mathcal{L}_{distill} = ...",
    "crop_path": ""
  },
  ...
]

Dependencies:
    pip install matplotlib python-pptx pillow

If you have a TeX distribution installed (TeX Live / MiKTeX), you can pass --usetex
for higher-quality rendering.
"""

import json,re
import argparse
from pathlib import Path
import matplotlib
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

import matplotlib
import matplotlib.pyplot as plt
from pathlib import Path
import shutil
HAS_TEX = shutil.which("latex") is not None  # 装了 TeX 才能用 usetex

def _render_once(latex: str, out_path: Path, dpi: int, fontsize: int, usetex: bool):
    
    try:
        matplotlib.rcParams["text.parse_math"] = True
    except Exception:
        pass
    matplotlib.rcParams.update({
        "text.usetex": (usetex and HAS_TEX),  # 未装 LaTeX 时为 False
        "mathtext.fontset": "cm",             # 可选
        "font.size": fontsize,
        "figure.dpi": dpi,
        "figure.max_open_warning": 0,         # 仅静音；真正释放靠 close(fig)
    })
     
    fig = plt.figure(figsize=(0.01, 0.01))
    try:
        ax = fig.add_subplot(111)
        ax.axis("off")
        txt = latex if (latex.startswith("$") and latex.endswith("$")) else f"${latex}$"
        ax.text(0.5, 0.5, txt, ha="center", va="center")
        out_path.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(out_path, dpi=dpi, bbox_inches="tight", pad_inches=0.12, transparent=True)
    finally:
        plt.close(fig)   

# 1) 先做基础空格压缩(避免命令与 { 之间/上下标周围的空格）
_SP_FIX = re.compile(
    r"(?:(?<=\\)\s+)|(?<=_)\s+|(?<=\^)\s+|(?<=\{)\s+|\s+(?=\})|\s+(?=[,=)])|(?<=\()\s+"
)

# \cmd { → \cmd{
_CMD_BRACE_FIX = re.compile(r"\\([A-Za-z]+)\s+\{")

def _collapse_spaces_inside_braces(text: str) -> str:
    # 把 { v i t } → {vit}，但保留逗号/下划线等分隔符
    def repl(m):
        inner = m.group(1)
        inner2 = re.sub(r"(?<=\w)\s+(?=\w)", "", inner)
        return "{" + inner2 + "}"
    return re.sub(r"\{([^{}]+)\}", repl, text)

def normalize_formula_strong(s: str) -> str:
    s = (s or "").strip()
    if not s:
        return s

    # 基础空格清理
    s = _SP_FIX.sub("", s)
    s = _CMD_BRACE_FIX.sub(r"\\\1{", s)
    s = _collapse_spaces_inside_braces(s)

    # 2) 修复常见 OCR/黏连错误 -------------------------

    # 1) \amalg 相关（把 "vm/llm" 识别成 \amalg m 的情况）
    #   - 先修下标里的 h_{\amalgm} 与 h_{\amalg m}
    s = re.sub(r"h_\{\\amalgm\}", r"h_{llm}", s)
    s = re.sub(r"h_\{\\amalg\s*m\}", r"h_{llm}", s)
    #   - 再兜底把 \amalgm / \amalg m → llm；孤立 \amalg → ll
    s = re.sub(r"\\amalgm", "llm", s)
    s = re.sub(r"\\amalg\s*m", "llm", s)
    s = re.sub(r"\\amalg\b", "ll", s)

    # 2) \mathrm 与 \vec 的非法嵌套/连写
    #   - 情况A：\mathrm\vec{vit} → \vec{vit}  （直接保留向量）
    s = re.sub(r"\\mathrm\\vec\{([^{}]+)\}", r"\\vec{\1}", s)
    #   - 情况B：\mathrm{\vec{vit}} → \vec{vit}
    s = re.sub(r"\\mathrm\{\\vec\{([^{}]+)\}\}", r"\\vec{\1}", s)
    #   - 情况C：\mathrmvit → \mathrm{vit}
    s = re.sub(r"\\mathrm([A-Za-z]+)", r"\\mathrm{\1}", s)
    #   - 情况D：\vecvit → \vec{vit}
    s = re.sub(r"\\vec([A-Za-z]+)", r"\\vec{\1}", s)

    # 3) \logP → \log P
    s = re.sub(r"\\log([A-Za-z])", r"\\log \1", s)

    # 4) 多个 \, 合并为一个
    s = re.sub(r"(\\,){2,}", r"\\,", s)

    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)
    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)

    # 2) \mathrmz / \mathrmd / \mathrmcp / \mathrmfux 等 → 加花括号（你已有单词版，这里补单字母）
    s = re.sub(r"\\mathrm([A-Za-z])", r"\\mathrm{\1}", s)
    s = re.sub(r"\\mathrm([A-Za-z])", r"\\mathrm{\1}", s)  # 再跑一遍兜底

    # 3) \mathbb{C} 等 → 用 \mathrm 代替（mathtext 对 \mathbb 支持有限）
    s = re.sub(r"\\mathbb\{([A-Za-z])\}", r"\\mathrm{\1}", s)

    # 4) 非法连写/伪命令
    s = s.replace(r"\logp",  r"\log p")
    s = s.replace(r"\circp", r"\circ p")
    s = s.replace(r"\cdotf", r"\cdot f")
    s = s.replace(r"\coloneqq", ":=")

    # 5) 奇怪的 \for 命令 → 用 \mathrm{for}（mathtext 不支持 \text）
    s = re.sub(r"\\for\b", r"\\mathrm{for}", s)

    # 6) 之前那条 \mathrm\vec{...} 组合再兜底一次
    s = re.sub(r"\\mathrm\\vec\{([^{}]+)\}", r"\\vec{\1}", s)
    s = re.sub(r"\\mathrm\{\\vec\{([^{}]+)\}\}", r"\\vec{\1}", s)

    # 7) 分布符号 \simD → \sim\mathcal{D}
    s = re.sub(r"\\simD\b", r"\\sim\\mathcal{D}", s)

    # (A) 单字母样式命令补花括号（\mathbfy → \mathbf{y}，\mathrmz → \mathrm{z}）
    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)
    s = re.sub(r"\\mathrm([A-Za-z])", r"\\mathrm{\1}", s)

    # (B) 黑板体在 mathtext 上不稳：\mathbb{C} → \mathrm{C}（可按需扩充）
    s = re.sub(r"\\mathbb\{([A-Za-z])\}", r"\\mathrm{\1}", s)

    # (C) 常见伪命令/连写
    s = s.replace(r"\logp",  r"\log p")
    s = s.replace(r"\circp", r"\circ p")
    s = s.replace(r"\cdotf", r"\cdot f")
    s = s.replace(r"\coloneqq", ":=")
    s = re.sub(r"\\for\b", r"\\mathrm{for}", s)   # \for → ‘for’

    # (D) 之前 VoRA 的 \amalg 误识别兜底
    s = re.sub(r"h_\{\\amalgm\}", r"h_{llm}", s)
    s = re.sub(r"\\amalgm",       "llm",      s)
    s = re.sub(r"\\amalg\s*m",    "llm",      s)
    s = re.sub(r"\\amalg\b",      "ll",       s)

    # (E) \mathrm\vec{vit}、\mathrm{\vec{vit}} → \vec{vit}
    s = re.sub(r"\\mathrm\\vec\{([^{}]+)\}",   r"\\vec{\1}", s)
    s = re.sub(r"\\mathrm\{\\vec\{([^{}]+)\}\}", r"\\vec{\1}", s)

    # (F) 二次下标：\nabla_{\mathrm{z}}_{t} → \nabla_{z_{t}}
    s = re.sub(r"\\nabla_\{\\mathrm\{([A-Za-z])\}\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{\\mathrm\{([A-Za-z])\}\}_([A-Za-z0-9]+)", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{([A-Za-z])\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)  # 更通用

    # (G) 分布符号：\simD / \simq → \sim \mathcal{D} / \sim q
    s = re.sub(r"\\simD\b", r"\\sim\\mathcal{D}", s)
    s = re.sub(r"\\simq\b", r"\\sim q", s)

    # (H) mathtext 不支持 \underbrace：\underbrace{A}_{B} → A
    s = re.sub(r"\\underbrace\{([^{}]+)\}\_\{[^{}]*\}", r"\1", s)

    # (I) 行尾截断：去掉尾部孤立 "_" 或 "\"，并补配对括号/花括号/left-right
    s = re.sub(r"[_\\]$", "", s)

    # 补齐 \left( / \right)
    n_left  = len(re.findall(r"\\left\(", s))
    n_right = len(re.findall(r"\\right\)", s))
    if n_left > n_right:
        s += "\\right)" * (n_left - n_right)

    # 补齐普通括号与花括号（粗暴但有效）
    diff_paren = s.count("(") - s.count(")")
    if diff_paren > 0:
        s += ")" * diff_paren
    diff_brace = s.count("{") - s.count("}")
    if diff_brace > 0:
        s += "}" * diff_brace
    # === 兼容 mathtext 的修复：覆盖你日志里的剩余错误 ===

    # 1) \simD / \simq：无论后面跟不跟下标都替换（不要用 \b）
    s = s.replace(r"\simD", r"\sim\mathcal{D}")
    s = s.replace(r"\simq", r"\sim q")

    # 2) \underbrace{A}_{B} → 仅保留 A（mathtext 不支持 underbrace）
    #   两条：先尽量匹配完整，再做一个宽松兜底
    s = re.sub(r"\\underbrace\s*\{([^{}]+)\}\s*_\s*\{[^{}]*\}", r"\1", s)
    s = re.sub(r"\\underbrace\s*\{([^{}]+)\}(\s*_\s*\{[^{}]*\})?", r"\1", s)

    # 3) 二次下标：\nabla_{\mathrm{z}}_{t} / \nabla_{z}_{t} → \nabla_{z_{t}}
    s = re.sub(r"\\nabla_\{\\mathrm\{([A-Za-z])\}\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{([A-Za-z])\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{([A-Za-z])\}_([A-Za-z0-9]+)", r"\\nabla_{\1_{\2}}", s)

    # 4) \mathbfX 这种 OCR 黏连（\mathbfx、\mathbfp、\mathbfV...）→ \mathbf{X}
    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)

    # 5) 行尾截断：去掉尾部残缺 token，并尽量补齐配对符号
    #   5.1 去掉收尾的“半个命令/下标/上标/反斜杠”
    s = re.sub(r"([_^]|\\)\s*$", "", s)

    #   5.2 补齐 \left( / \left[ / \left\{ / \left\| 与对应的 \right
    def _balance(pair_left, pair_right):
        nL = len(re.findall(pair_left, s))
        nR = len(re.findall(pair_right, s))
        return nL - nR
    diff = _balance(r"\\left\(", r"\\right\)")
    if diff > 0: s += "\\right)" * diff
    diff = _balance(r"\\left\[", r"\\right\]")
    if diff > 0: s += "\\right]" * diff
    diff = _balance(r"\\left\\\{", r"\\right\\\}")
    if diff > 0: s += "\\right\\}" * diff
    diff = _balance(r"\\left\\\|", r"\\right\\\|")
    if diff > 0: s += "\\right\\|" * diff

    #   5.3 普通括号/花括号配对（粗暴但有效）
    diff_paren = s.count("(") - s.count(")")
    if diff_paren > 0: s += ")" * diff_paren
    diff_brace = s.count("{") - s.count("}")
    if diff_brace > 0: s += "}" * diff_brace

    # 连续空白再压一次
    s = re.sub(r"\s+", " ", s).strip()
   
    return s
def render_latex_to_png(latex: str, out_path: Path, dpi: int = 220, fontsize: int = 24, usetex: bool = False):
    # 先做强力规范化
    latex_norm = normalize_formula_strong(latex)
    try:
        _render_once(latex_norm, out_path, dpi, fontsize, usetex=False)
        return
    except Exception as e1:
        if usetex:
            # 用户显式要求 usetex
            _render_once(latex_norm, out_path, dpi, fontsize, usetex=True)
            return
        else:
            # mathtext 失败，自动兜底 usetex
            _render_once(latex_norm, out_path, dpi, fontsize, usetex=True)
            return



def formulas_json_to_ppt(
    json_path: Path,
    ppt_path: Path,
    img_dir: Path,
    dpi: int = 220,
    fontsize: int = 24,
    usetex: bool = False,
    max_per_slide: int = 1,
    margin_in: float = 1.0,
):
    """
    Read formulas from JSON; render each as PNG; insert into a PPT.

    Args:
        json_path: path to JSON file
        ppt_path: output PPTX path
        img_dir: directory to store rendered images
        dpi: PNG dpi
        fontsize: LaTeX font size (points)
        usetex: use system LaTeX (needs TeX Live/MiKTeX); fallback is mathtext
        max_per_slide: how many formulas per slide (1 = one per slide)
        margin_in: margin (inches) from slide borders
    """
    data = json.loads(Path(json_path).read_text(encoding="utf-8"))

    # Accept either a list or an object with "items"
    if isinstance(data, dict) and "items" in data:
        items = data["items"]
    else:
        items = data

    # Build list of latex strings
    formulas = []
    for it in items:
        latex = it.get("latex") or it.get("latex_raw") or ""
        latex = latex.strip()
        if latex:
            formulas.append(latex)

    if not formulas:
        print("No formulas found in JSON (looked for 'latex' or 'latex_raw').")
        return

    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Convert EMU to inches helpers (pptx stores sizes in EMUs)
    EMU_PER_INCH = 914400
    slide_w_in = slide_w / EMU_PER_INCH
    slide_h_in = slide_h / EMU_PER_INCH

    img_dir.mkdir(parents=True, exist_ok=True)

    # Render and place
    slide = None
    per_slide_count = 0

    for idx, latex in enumerate(formulas, start=1):
        # 1) Render PNG
        img_path = img_dir / f"formula_{idx:04d}.png"
        try:
            print("input :", latex)
            render_latex_to_png(latex, img_path, dpi=dpi, fontsize=fontsize, usetex=usetex)
        except Exception as e:
            print(f"[Warning] Failed to render formula #{idx}: {e}")
            continue

        # 2) Add slide when needed
        if slide is None or per_slide_count >= max_per_slide:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            per_slide_count = 0

        # 3) Compute placement
        # Try to keep image within slide, centered; width = slide_w - 2*margin
        content_w_in = max(0.5, slide_w_in - 2 * margin_in)
        left_in = margin_in
        # default top: center vertically if one per slide; otherwise stack
        if max_per_slide == 1:
            top_in = (slide_h_in - content_w_in * 0.5) / 2  # rough center guess
            if top_in < margin_in:
                top_in = margin_in
        else:
            # stack with equal spacing
            row_h_in = (slide_h_in - 2 * margin_in) / max_per_slide
            top_in = margin_in + per_slide_count * row_h_in + (row_h_in - row_h_in * 0.6) / 2
            content_w_in = content_w_in  # keep width

        # Insert picture; pptx will scale height to preserve aspect ratio if only width given
        pic = slide.shapes.add_picture(
            str(img_path),
            left=int(left_in * EMU_PER_INCH),
            top=int(top_in * EMU_PER_INCH),
            width=int(content_w_in * EMU_PER_INCH),
            height=None,
        )

        # Optional: add a small caption text placeholder (commented)
        # tx = slide.shapes.add_textbox(
        #     left=int(margin_in * EMU_PER_INCH),
        #     top=int((top_in + (pic.height/EMU_PER_INCH) + 0.2) * EMU_PER_INCH),
        #     width=int((slide_w_in - 2*margin_in) * EMU_PER_INCH),
        #     height=int(0.5 * EMU_PER_INCH),
        # )
        # tx.text_frame.text = f"Formula {idx}"

        per_slide_count += 1

    prs.save(ppt_path)
    print(f"Done. Rendered {len(formulas)} formulas.")
    print(f"PPT saved to: {ppt_path}")
    print(f"PNGs saved under: {img_dir}")

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--json", required=True, help="Path to formulas JSON (array of items)")
    ap.add_argument("--out-ppt", default="formulas.pptx", help="Output PPTX path")
    ap.add_argument("--img-dir", default="formula_imgs", help="Directory to save rendered PNGs")
    ap.add_argument("--dpi", type=int, default=220, help="PNG DPI (default 220)")
    ap.add_argument("--fontsize", type=int, default=24, help="LaTeX font size (pt)")
    ap.add_argument("--usetex", action="store_true", help="Use system LaTeX (needs TeX Live/MiKTeX)")
    ap.add_argument("--max-per-slide", type=int, default=1, help="Formulas per slide (default 1)")
    ap.add_argument("--margin-in", type=float, default=1.0, help="Slide margin in inches (default 1.0)")
    args = ap.parse_args()

    formulas_json_to_ppt(
        json_path=Path(args.json),
        ppt_path=Path(args.out_ppt),
        img_dir=Path(args.img_dir),
        dpi=args.dpi,
        fontsize=args.fontsize,
        usetex=args.usetex,
        max_per_slide=args.max_per_slide,
        margin_in=args.margin_in,
    )

if __name__ == "__main__":
    main()



'''
# 基本用法(不依赖 TeX 安装，使用 matplotlib 内置 mathtext）
python formulas_json_to_ppt.py \
  --json formulas_with_bbox.json \
  --out-ppt formulas.pptx \
  --img-dir formula_imgs \
  --dpi 220 \
  --fontsize 26

# 如果你机器装了 TeX(TeX Live / MiKTeX 等），可加 --usetex 提升渲染质量
python formulas_json_to_ppt.py \
  --json formulas_with_bbox.json \
  --out-ppt formulas_usetex.pptx \
  --img-dir formula_imgs_tex \
  --dpi 250 \
  --fontsize 26 \
  --usetex





'''